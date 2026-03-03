# =================================================================
# 1. CORE CONFIGURATION & IMPORTS
# =================================================================
from flask import Flask, request, jsonify, render_template
import chromadb, json, datetime
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from groq import Groq
from dotenv import load_dotenv
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image, ImageOps, ImageEnhance
import os, uuid, io, re, socket
from concurrent.futures import ThreadPoolExecutor
from functools import partial

pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_PATH', r'C:\Program Files\Tesseract-OCR\tesseract.exe')

# Tesseract configuration for better OCR and handwritten notes
# --oem 3: Default, based on what is available.
# --psm 11: Sparse text. Find as much text as possible in no particular order.
TESSERACT_CONFIG = r'--oem 3 --psm 11'

def clean_ocr_text(text):
    """Filter out common OCR noise and scanner watermarks WITHOUT deleting real content"""
    if not text:
        return ""
        
    watermarks = [
        r"OKEN Scanner",
        r"CamScanner",
        r"CS CamScanner",
        r"Scanned with",
        r"Shot on",
        r"Watermark"
    ]
    
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        line_strip = line.strip()
        if not line_strip:
            continue
            
        is_watermark = False
        # If line contains both 'OKEN' and 'Scanner', or 'Cam' and 'Scanner', it's likely a watermark
        # Even if it has noise around it like "@ Scanned with OKEN Scanner ."
        lower_line = line_strip.lower()
        
        # Check for specific combinations that are highly likely to be watermarks
        if ("oken" in lower_line and "scanner" in lower_line) or \
           ("cam" in lower_line and "scanner" in lower_line) or \
           ("scanned" in lower_line and "with" in lower_line and len(line_strip) < 40):
            is_watermark = True
        
        # Additional check for short lines matching any watermark term
        if not is_watermark and len(line_strip) < 30:
            for wm in watermarks:
                if wm.lower() in lower_line:
                    is_watermark = True
                    break
        
        if not is_watermark:
            cleaned_lines.append(line)
    
    # Rejoin and remove excessive newlines
    result = '\n'.join(cleaned_lines)
    result = re.sub(r'\n\s*\n', '\n', result)
    return result.strip()

# Load environment variables
load_dotenv()

def get_ip():
    """Returns local network IP address"""
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(('8.8.8.8', 80))
        IP = s.getsockname()[0]
        s.close()
        if IP and IP != '127.0.0.1':
            return IP
    except Exception:
        pass
    
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(('1.1.1.1', 53))
        IP = s.getsockname()[0]
        s.close()
        if IP and IP != '127.0.0.1':
            return IP
    except Exception:
        pass
    
    try:
        hostname = socket.gethostname()
        IP = socket.gethostbyname(hostname)
        if IP and IP != '127.0.0.1':
            return IP
    except Exception:
        pass
    
    return '127.0.0.1'

os.environ["TF_CPP_MIN_LOG_LEVEL"] = "3"

app = Flask(__name__)

# Chat history storage
CHATS_FILE = os.path.join(os.getcwd(), "chats.json")
if not os.path.exists(CHATS_FILE):
    with open(CHATS_FILE, "w") as f:
        json.dump({}, f)

def get_chats():
    with open(CHATS_FILE, "r") as f:
        return json.load(f)

def save_chats(chats):
    with open(CHATS_FILE, "w") as f:
        json.dump(chats, f)

# Groq Client Initialization
groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY"))
MODEL_NAME = "llama-3.3-70b-versatile"

def call_groq(messages):
    """Utility to call Groq LLM API"""
    try:
        completion = groq_client.chat.completions.create(
            model=MODEL_NAME,
            messages=messages,
        )
        return completion.choices[0].message.content
    except Exception as e:
        raise Exception(f"Groq API Error: {str(e)}")

# =================================================================
# 2. VECTOR DATABASE (ChromaDB) SETUP
# =================================================================
db_path = os.path.join(os.getcwd(), "chroma_db")
client = chromadb.PersistentClient(path=db_path)
collection = client.get_or_create_collection("documents")

# Embedding Model Initialization
embedder = SentenceTransformer("all-MiniLM-L6-v2")

# =================================================================
# 3. TEXT EXTRACTION & CHUNKING
# =================================================================
def extract_text(file_stream, filename):
    """Extract raw text from PDF, DOCX, PPTX, and TXT files"""
    fn_lower = filename.lower()
    if fn_lower.endswith(".pdf"):
        try:
            file_stream.seek(0)
            text_parts = []
            cleaned_initial_text = ""
            
            try:
                reader = PdfReader(file_stream)
                # Check for encryption
                if reader.is_encrypted:
                    try:
                        reader.decrypt("")
                    except Exception:
                        pass # Try anyway or rely on OCR
                
                for page_idx, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text_parts.append(page_text)
                    
                    try:
                        if "/XObject" in page["/Resources"]:
                            xObject = page["/Resources"]["/XObject"].get_object()
                            for obj in xObject:
                                if xObject[obj]["/Subtype"] == "/Image":
                                    try:
                                        from PIL import Image as PILImage
                                        size = (int(xObject[obj]["/Width"]), int(xObject[obj]["/Height"]))
                                        data = xObject[obj].get_data()
                                        image = PILImage.frombytes("RGB", size, data)
                                        
                                        img_text = pytesseract.image_to_string(image, config=TESSERACT_CONFIG)
                                        img_text = clean_ocr_text(img_text)
                                        if img_text.strip():
                                            text_parts.append(f"[IMAGE CONTENT]: {img_text}")
                                        else:
                                            text_parts.append(f"[EMPTY IMAGE]")
                                    except Exception:
                                        pass
                    except Exception:
                        pass
                
                text = "\n".join(text_parts)
                cleaned_initial_text = clean_ocr_text(text)
            except Exception as e:
                print(f"DEBUG: pypdf failed to parse {filename}: {e}. Falling back to full OCR.")
                cleaned_initial_text = ""
            
            # OCR fallback if text is sparse or pypdf failed
            if len(cleaned_initial_text.strip()) < 200:
                try:
                    print(f"DEBUG: Attempting OCR for {filename} (Cleaned length: {len(cleaned_initial_text)})")
                    file_stream.seek(0)
                    file_bytes = file_stream.read()
                    if not file_bytes:
                        return "ERROR: PDF file is empty."
                        
                    poppler_path = os.getenv('POPPLER_PATH', r'C:\poppler\poppler-24.08.0\Library\bin')
                    # Optimize: Lower DPI (150) and use multiple threads for conversion
                    images = convert_from_bytes(
                        file_bytes, 
                        poppler_path=poppler_path, 
                        dpi=300, # Higher DPI for better handwriting recognition
                        fmt="jpeg", 
                        thread_count=os.cpu_count() or 4
                    )
                    
                    if not images:
                        return "ERROR: Could not convert PDF to images for OCR. Check Poppler installation."
                        
                    print(f"DEBUG: Converted PDF to {len(images)} images at 300 DPI")
                    
                    # Optimize: Parallel OCR processing
                    def perform_ocr(image, index):
                        try:
                            print(f"DEBUG: Starting OCR on page {index+1}")
                            # Image Pre-processing for better handwritten note visibility
                            img = ImageOps.grayscale(image)
                            img = ImageOps.autocontrast(img)
                            
                            # Boost contrast significantly for handwriting
                            contrast = ImageEnhance.Contrast(img)
                            img = contrast.enhance(2.0)
                            
                            # Sharpen to make lines crisper
                            enhancer = ImageEnhance.Sharpness(img)
                            img = enhancer.enhance(2.0)
                            
                            text = pytesseract.image_to_string(img, config=TESSERACT_CONFIG)
                            return index, text
                        except Exception as e:
                            print(f"DEBUG: OCR thread error for page {index+1}: {e}")
                            return index, ""

                    ocr_results_map = {}
                    with ThreadPoolExecutor(max_workers=os.cpu_count() or 4) as executor:
                        futures = [executor.submit(perform_ocr, img, i) for i, img in enumerate(images)]
                        for future in futures:
                            idx, page_ocr = future.result()
                            if page_ocr.strip():
                                ocr_results_map[idx] = page_ocr
                    
                    ocr_results = [ocr_results_map[idx] for idx in sorted(ocr_results_map.keys())]
                    ocr_text_raw = "\n".join(ocr_results)
                    ocr_text = clean_ocr_text(ocr_text_raw)
                    
                    # LOGGING: See what we actually got after cleaning
                    print(f"DEBUG: OCR extracted {len(ocr_text_raw)} raw characters, {len(ocr_text)} after cleaning")
                    
                    if len(ocr_text.strip()) > len(cleaned_initial_text.strip()):
                        text = ocr_text
                        print(f"DEBUG: OCR successfully extracted {len(text)} characters")
                    else:
                        text = cleaned_initial_text
                        print("DEBUG: OCR extracted less text than original method")
                except Exception as e:
                    print(f"DEBUG: OCR Error: {str(e)}")
                    if "poppler" in str(e).lower():
                        return f"ERROR: Poppler not found at {poppler_path}. Please ensure Poppler is extracted there."
                    return f"ERROR: OCR failed - {str(e)}"
            
            else:
                text = cleaned_initial_text
            
            return text
        except Exception as e:
            print(f"DEBUG: PDF Reader Error: {str(e)}")
            return f"ERROR: PDF extraction failed - {str(e)}"
    
    if fn_lower.endswith(".docx"):
        doc = Document(file_stream)
        text = [p.text for p in doc.paragraphs]
        
        for table_idx, table in enumerate(doc.tables):
            text.append(f"\n[TABLE {table_idx + 1}]")
            for row_idx, row in enumerate(table.rows):
                row_cells = [cell.text.strip() for cell in row.cells]
                text.append(" | ".join(row_cells))
        
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    image_part = rel.target_part
                    image_data = image_part.blob
                    image = Image.open(io.BytesIO(image_data))
                    img_text = pytesseract.image_to_string(image, config=TESSERACT_CONFIG)
                    img_text = clean_ocr_text(img_text)
                    if img_text.strip():
                        text.append(f"\n[IMAGE CONTENT]: {img_text}")
                except Exception:
                    pass
        
        return "\n".join(text)
    
    if fn_lower.endswith(".pptx") or fn_lower.endswith(".ppt"):
        try:
            file_stream.seek(0)
            stream = io.BytesIO(file_stream.read())
            prs = Presentation(stream)
            text = []
            for slide_idx, slide in enumerate(prs.slides):
                text.append(f"\n[SLIDE {slide_idx + 1}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
                    if shape.has_table:
                        table = shape.table
                        text.append("[TABLE]")
                        for row in table.rows:
                            row_cells = [cell.text.strip() for cell in row.cells]
                            text.append(" | ".join(row_cells))
                    if shape.shape_type == 13:
                        try:
                            image_stream = io.BytesIO(shape.image.blob)
                            image = Image.open(image_stream)
                            img_text = pytesseract.image_to_string(image, config=TESSERACT_CONFIG)
                            img_text = clean_ocr_text(img_text)
                            if img_text.strip():
                                text.append(f"[IMAGE]: {img_text}")
                        except Exception:
                            pass
            return "\n".join(text)
        except Exception:
            if fn_lower.endswith(".ppt"):
                return "ERROR: Legacy .ppt detected. Please save as .pptx."
            return "ERROR: Could not read PowerPoint file."
    
    if fn_lower.endswith(".txt"):
        return file_stream.read().decode("utf-8")
    
    image_extensions = [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff", ".webp"]
    if any(fn_lower.endswith(ext) for ext in image_extensions):
        try:
            file_stream.seek(0)
            image = Image.open(file_stream)
            raw_text = pytesseract.image_to_string(image, config=TESSERACT_CONFIG)
            return clean_ocr_text(raw_text)
        except Exception:
            return "ERROR: Could not read image file."
    
    return ""

def chunk_text(text, chunk_size=500, overlap=50):
    """Split text into overlapping chunks for indexing"""
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks

def ingest_file(file_stream, filename):
    """Process and store file in vector database"""
    text = extract_text(file_stream, filename)
    if text.startswith("ERROR:"):
        return text
    if not text.strip():
        return 0
        
    chunks = chunk_text(text)
    embeddings = embedder.encode(chunks, normalize_embeddings=True).tolist()
    
    file_type = filename.split('.')[-1].lower()
    ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
    metadatas = [{"source": filename, "chunk_id": idx, "type": file_type} for idx in range(len(chunks))]
    
    collection.add(
        documents=chunks,
        embeddings=embeddings,
        metadatas=metadatas,
        ids=ids
    )
    return len(chunks)

# =================================================================
# 4. RETRIEVAL & CONTEXT GENERATION
# =================================================================
def retrieve(query, where=None, n_results=5):
    """Search for relevant document chunks"""
    q_emb = embedder.encode([query], normalize_embeddings=True).tolist()
    res = collection.query(
        query_embeddings=q_emb,
        n_results=n_results,
        include=["documents", "metadatas", "distances"],
        where=where
    )
    
    docs = []
    metas = []
    if res["documents"] and len(res["documents"]) > 0:
        for d, m, dist in zip(res["documents"][0], res["metadatas"][0], res["distances"][0]):
            if dist < 1.6: 
                docs.append(d)
                metas.append(m)
    return docs, metas

def generate_standalone_question(question, history):
    """Rephrase user question based on conversation history"""
    if not history:
        return question
    
    history_str = "\n".join([f"{m['role'].upper()}: {m['content']}" for m in history[-5:]])
    prompt = f"""Conversation History:
{history_str}

Follow-up Question: {question}
Rephrase the follow-up question to be a standalone search query. Only return the query."""
    
    return call_groq([{"role": "user", "content": prompt}]).strip()

# =================================================================
# 5. FLASK ROUTES
# =================================================================
@app.route("/")
def home():
    return render_template("index.html", title="DOCUMENT SUMMARIZER", ip="DOCUMENT_SUMMARIZER.local:8080")

@app.route("/upload", methods=["POST"])
def upload():
    """Handle document uploads"""
    try:
        results = []
        for f in request.files.getlist("files"):
            result = ingest_file(f, f.filename)
            if isinstance(result, str) and result.startswith("ERROR:"):
                results.append({"filename": f.filename, "error": result.replace("ERROR:", "").strip()})
            else:
                results.append({"filename": f.filename, "chunks": result})
        return jsonify({"status": "Upload complete", "files": results})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/ask", methods=["POST"])
def ask():
    """Handle chat queries using RAG"""
    try:
        data = request.json
        q = data["question"]
        history = data.get("history", [])
        
        # Determine filters based on keywords
        where_filter = None
        target_type = None
        q_lower = q.lower()
        if "ppt" in q_lower or "powerpoint" in q_lower:
            target_type = "pptx"
            where_filter = {"type": {"$in": ["pptx", "ppt"]}}
        elif "pdf" in q_lower:
            target_type = "pdf"
            where_filter = {"type": "pdf"}
        elif "word" in q_lower or "docx" in q_lower:
            target_type = "docx"
            where_filter = {"type": "docx"}
        elif "text" in q_lower or "txt" in q_lower:
            target_type = "txt"
            where_filter = {"type": "txt"}

        # Handle broad queries
        general_queries = ["what is this", "summary of", "about this", "describe", "what is the"]
        is_general = any(x in q_lower for x in general_queries)
        
        if is_general:
            if target_type:
                where_clause = {"type": target_type} if isinstance(target_type, str) else where_filter
                data_files = collection.get(where=where_clause, include=["metadatas"])
            else:
                data_files = collection.get(include=["metadatas"])
                
            filenames = list(set(m["source"] for m in data_files["metadatas"]))
            
            if len(filenames) > 1:
                mentioned = [f for f in filenames if f.lower() in q_lower]
                if len(mentioned) == 1:
                    where_filter = {"source": mentioned[0]}
                elif not mentioned:
                    return jsonify({"answer": f"Please specify which file you mean: {', '.join(filenames)}"})
            elif len(filenames) == 1:
                where_filter = {"source": filenames[0]}

        # RAG Pipeline: Rephrase -> Retrieve -> Generate
        standalone_q = generate_standalone_question(q, history)
        n_results = 15 if is_general else 5
        docs, metas = retrieve(standalone_q, where=where_filter, n_results=n_results)
        context = "\n".join(docs) if docs else "No relevant context found."

        system_prompt = f"""You are the DOCUMENT SUMMARIZER. Use the provided context to answer.
1. If greeting, respond politely.
2. If info is in context, answer in detail.
3. The context may contain OCR noise, errors, or scanner watermarks (like "OKEN Scanner", "CamScanner", etc.). 
   - **IMPORTANT**: Ignore all scanner watermarks and boilerplate text.
4. **IMAGE CAPABILITIES**:
   - **AI DIAGRAMS**: If asked for a diagram, flowchart, or graph, generate it using Mermaid.js code blocks: ```mermaid [code] ```. **IMPORTANT**: Always start the mermaid code with `graph TD` or `flowchart TD` to ensure it renders correctly. Avoid using parentheses `()` or brackets `[]` inside node names unless necessary, as they can cause syntax errors.
   - **AI IMAGE GENERATION**: If the user asks to "generate an image", "show a picture of", or if a visual would help explain a concept, use Pollinations.ai: `![AI Image](https://image.pollinations.ai/prompt/description)`. Replace 'description' with a URL-safe description.
5. Use your intelligence to correct obvious typos and infer meaning from handwritten notes.
6. If info is NOT in context, respond with EXACTLY: "This information is not available in the provided documents."
7. End response with: "Sources used: filename1, filename2". If none, "Sources used: None".

Context:
{context}"""
        
        formatted_history = []
        for msg in history[-10:]:
            role = msg.get("role", "user")
            if role == "ai":
                role = "assistant"
            elif role not in ["user", "assistant"]:
                role = "user"
            formatted_history.append({"role": role, "content": msg.get("content", "")})
        
        messages = [{"role": "system", "content": system_prompt}] + formatted_history + [{"role": "user", "content": q}]
        answer = call_groq(messages)

        # Save to chat history if session_id is provided
        session_id = data.get("session_id")
        if session_id:
            chats = get_chats()
            if session_id not in chats:
                chats[session_id] = {
                    "messages": [],
                    "title": "New Chat",
                    "timestamp": datetime.datetime.now().isoformat()
                }
            
            chats[session_id]["messages"].append({"role": "user", "content": q})
            chats[session_id]["messages"].append({"role": "ai", "content": answer, "sources": metas})
            
            # Update title if it's the first message
            if chats[session_id]["title"] == "New Chat":
                # Generate an apt title using LLM
                try:
                    title_prompt = f"Generate a very short, descriptive title (max 5 words) for a conversation that starts with this question: '{q}'. Respond with ONLY the title text."
                    title = call_groq([{"role": "user", "content": title_prompt}]).strip().replace('"', '').replace("'", "")
                except:
                    # Fallback to simple logic if LLM title fails
                    title = " ".join(q.split()[:5])
                    if len(q.split()) > 5:
                        title += "..."
                chats[session_id]["title"] = title
            
            save_chats(chats)

        # Cleanup citations and sources
        used_sources = []
        if "Sources used:" in answer:
            parts = answer.split("Sources used:")
            source_line = parts[-1].strip()
            answer = parts[0].strip()
            used_sources = [s.strip() for s in source_line.split(",")]
            
        patterns = [r"(?i)\n\s*File names:.*", r"(?i)\n\s*Sources:.*", r"(?i)\n\s*Files:.*", r"(?i)Sources used:.*"]
        for p in patterns:
            answer = re.split(p, answer)[0].strip()

        final_sources = []
        seen = set()
        for m in metas:
            # Only include if the source was actually mentioned by the AI
            source_file = m["source"]
            is_used = any(source_file.lower() in s.lower() for s in used_sources) if used_sources else False
            
            if is_used or (not used_sources and context != "No relevant context found."):
                if (source_file, m["chunk_id"]) not in seen:
                    final_sources.append({"file": source_file, "chunk": m["chunk_id"]})
                    seen.add((source_file, m["chunk_id"]))

        return jsonify({"answer": answer, "sources": final_sources})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/summary", methods=["POST"])
def summary():
    """Generate comprehensive document summaries"""
    try:
        filename = request.json.get("filename")
        where = {"source": filename} if filename else None
        data = collection.get(where=where, include=["documents", "metadatas"])
            
        if not data["documents"]:
            return jsonify({"error": "No content found"}), 404
            
        combined = sorted(zip(data["metadatas"], data["documents"]), key=lambda x: x[0].get("chunk_id", 0))
        text = "\n".join([x[1] for x in combined[:20]])
        summary_text = call_groq([
            {"role": "system", "content": "You are a professional document summarizer. Create a concise yet comprehensive summary of the provided text with key highlights and main points. Use Markdown for formatting."},
            {"role": "user", "content": f"Please summarize the following document content:\n\n{text}"}
        ])

        return jsonify({
            "summary": summary_text,
            "sources": [{"file": m["source"], "chunk": m["chunk_id"]} for m in data["metadatas"]]
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/clear", methods=["POST"])
def clear_db():
    """Reset the vector database and chat history"""
    try:
        global collection
        client.delete_collection("documents")
        collection = client.create_collection("documents")
        
        # Clear chat history
        save_chats({})
        
        return jsonify({"status": "Database and history cleared"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/files")
def list_files():
    """List all indexed files"""
    try:
        data = collection.get(include=["metadatas"])
        sources = sorted(list(set(m["source"] for m in data["metadatas"])))
        return jsonify({"files": sources})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/chats", methods=["GET"])
def list_chats():
    """List all chat sessions"""
    try:
        chats = get_chats()
        history = []
        for chat_id, data in chats.items():
            # Only list chats that have messages
            if data.get("messages"):
                title = data.get("title", "New Chat")
                history.append({"id": chat_id, "title": title, "timestamp": data.get("timestamp")})
        return jsonify({"chats": history})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/chat/<chat_id>", methods=["GET"])
def get_chat(chat_id):
    """Retrieve messages for a specific chat"""
    try:
        chats = get_chats()
        if chat_id in chats:
            return jsonify(chats[chat_id])
        return jsonify({"error": "Chat not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/delete_chat/<chat_id>", methods=["DELETE"])
def delete_chat(chat_id):
    """Remove a specific chat session"""
    try:
        chats = get_chats()
        if chat_id in chats:
            del chats[chat_id]
            save_chats(chats)
            return jsonify({"status": "Chat deleted"})
        return jsonify({"error": "Chat not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/rename_chat/<chat_id>", methods=["POST"])
def rename_chat(chat_id):
    """Change the title of a chat session"""
    try:
        data = request.json
        new_title = data.get("title")
        if not new_title:
            return jsonify({"error": "Title is required"}), 400
            
        chats = get_chats()
        if chat_id in chats:
            chats[chat_id]["title"] = new_title
            save_chats(chats)
            return jsonify({"status": "Chat renamed"})
        return jsonify({"error": "Chat not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/new_chat", methods=["POST"])
def new_chat():
    """Return a new session ID for a chat"""
    try:
        chat_id = str(uuid.uuid4())
        return jsonify({"chat_id": chat_id})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# =================================================================
# 6. APP RUNNER
# =================================================================
if __name__ == "__main__":
    app.run(debug=True)
